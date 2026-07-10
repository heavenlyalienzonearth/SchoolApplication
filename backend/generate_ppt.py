import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # Color Palette (Kid-friendly Professional Style)
    COLOR_BG = RGBColor(255, 248, 245)       # Warm off-white
    COLOR_PRIMARY = RGBColor(30, 41, 59)      # Slate Dark Navy
    COLOR_SECONDARY = RGBColor(238, 90, 36)    # Orange
    COLOR_ACCENT = RGBColor(2, 132, 199)      # Sky Blue
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT = RGBColor(71, 85, 105)        # Muted Text

    # Helper: Set Slide Background
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Add Standard Header
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.name = 'Helvetica'
        
        # Bottom underline bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.3), Inches(2.0), Inches(0.06))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SECONDARY
        shape.line.fill.background()

    # Helper: Add Card
    def add_card(slide, left, top, width, height, title, body_bullets):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.line.color.rgb = RGBColor(226, 232, 240)
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_bottom = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.name = 'Helvetica'
        p.alignment = PP_ALIGN.LEFT
        
        for bullet in body_bullets:
            p2 = tf.add_paragraph()
            p2.text = "• " + bullet
            p2.font.size = Pt(13)
            p2.font.color.rgb = COLOR_TEXT
            p2.font.name = 'Helvetica'
            p2.space_before = Pt(6)
            p2.alignment = PP_ALIGN.LEFT

    # Helper: Add Slide Layout with Card & Image
    def add_slide_layout(slide, header_title, card_title, bullets, image_filename):
        # Add Slide Header
        add_slide_header(slide, header_title)
        
        # Add descriptive card on the left
        add_card(slide, Inches(0.75), Inches(1.8), Inches(5.2), Inches(4.8), card_title, bullets)
        
        # Add screenshot on the right
        image_path = os.path.join("screenshots", image_filename)
        if os.path.exists(image_path):
            # Frame background card
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(1.8), Inches(6.23), Inches(4.8))
            frame.fill.solid()
            frame.fill.fore_color.rgb = COLOR_WHITE
            frame.line.color.rgb = RGBColor(226, 232, 240)
            frame.line.width = Pt(1.5)
            
            # Picture
            slide.shapes.add_picture(image_path, Inches(6.45), Inches(1.9), width=Inches(6.03), height=Inches(4.6))
        else:
            # Placeholder if image missing
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.35), Inches(1.8), Inches(6.23), Inches(4.8))
            frame.fill.solid()
            frame.fill.fore_color.rgb = COLOR_WHITE
            frame.line.color.rgb = COLOR_SECONDARY
            
            tf = frame.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"📷 UI Screenshot Placeholder\n({image_filename})"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT
            p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, COLOR_SECONDARY)

    # Left colored accent card
    left_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = COLOR_PRIMARY
    left_accent.line.fill.background()

    # Title box
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "Kangaroo Kids Portal"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = 'Helvetica'

    p_sub = tf1.add_paragraph()
    p_sub.text = "Comprehensive System Walkthrough & Feature Overview"
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = COLOR_PRIMARY
    p_sub.font.name = 'Helvetica'
    p_sub.space_before = Pt(10)

    # Footer/Metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(6.0), Inches(1.0))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Created: July 2026\nTarget Audience: School Administration & Stakeholders"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = COLOR_WHITE
    p_meta.font.italic = True
    p_meta.font.name = 'Helvetica'

    # ==========================================
    # SLIDE 2: SITE CONFIGURATION
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, COLOR_BG)
    add_slide_layout(
        slide2,
        "⚙️ Site Settings & Global Branding",
        "Branding & Configuration",
        [
            "School Logo Upload: Supports editing and configuring the global logo.",
            "Name & Branding: Configurable school name synced across all parent headers.",
            "Contact Information: Dynamic address, phone, and official support email entries.",
            "Government Approval Details: Added dedicated approval input fields for Karnataka education regulations.",
            "Dynamic Binding: Automatically inserts registration codes on promotion cards and leaving forms."
        ],
        "dashboard_analytics.png"
    )

    # ==========================================
    # SLIDE 3: HOMEPAGE & EXPERIENCE
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3, COLOR_BG)
    add_slide_layout(
        slide3,
        "🖼️ Homepage & Parent Experience",
        "Public Interface Features",
        [
            "Dynamic Hero Slider: Engaging custom banners highlighting preschool events.",
            "Learning Programs: Curated cards describing Toddler, Preschool, and Kindergarten curricula.",
            "Parent Reviews Slider: Visual, animated slider displaying positive parent testimonials.",
            "School Photo Gallery: Slide preview showcasing creative playtime activities.",
            "Virtual Assistant Chatbot: Smart support chatbot answering visitor questions regarding admissions, fees, and timings."
        ],
        "homepage.png"
    )

    # ==========================================
    # SLIDE 4: ADMISSIONS SYSTEM
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4, COLOR_BG)
    add_slide_layout(
        slide4,
        "🏫 Admissions & Vaccinations Manager",
        "Student Admissions Portal",
        [
            "Digital Onboarding Form: Simple parent inputs for child, parent, email, and phone.",
            "Age-Group Admissions: Automated grade-specific enrollment paths.",
            "Uniform Checklist: Built-in selections for issued uniforms, bags, caps, and shoes.",
            "Contact Records: Collects emergency phone contacts and student blood groups.",
            "Medical Safety Tracker: Integrates immunization schedules with registration."
        ],
        "dashboard_admissions.png"
    )

    # ==========================================
    # SLIDE 5: ATTENDANCE & ROSTER
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5, COLOR_BG)
    add_slide_layout(
        slide5,
        "📅 Daily Attendance & Class Roster",
        "Attendance marking & Class Roster",
        [
            "Class Selector: Easy dropdown filter to pick classes (e.g. Preschool Engineers).",
            "Student Checklist: Generates a list of students with photo thumbnails.",
            "Status Toggle: Mark Present or Absent with one click.",
            "Quick Save: Real-time update to the attendance log database.",
            "Enrolled Class Roster: Highlights student details, issued items, and interactive certificate controls."
        ],
        "transfer_certificate.png" # Standard list layout screenshot
    )

    # ==========================================
    # SLIDE 6: ID CARD GENERATION
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6, COLOR_BG)
    add_slide_layout(
        slide6,
        "📇 Automated Student ID Badges",
        "ID Card Design & Mailing",
        [
            "Widescreen Double-Sided Badge: Renders both Front and Back cards side-by-side.",
            "Front Details: Child Name, Photo, Grade, Class Current Year (e.g. 2026), Blood Group, and Emergency Phone.",
            "Privacy Focus: Removed parent names from the badge front for student safety.",
            "Scannable Barcode: Embeds a scannable Libre Barcode 39 font linked to student ID.",
            "One-Click Mailing: Send button formats and requests badges directly from printshops."
        ],
        "id_badge_modal.png"
    )

    # ==========================================
    # SLIDE 7: OFFICIAL CERTIFICATES
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7, COLOR_BG)
    add_slide_layout(
        slide7,
        "📜 Govt. Certificates & TC Printing",
        "Promotion & Transfer Cards",
        [
            "Department Compliance: Formal certificate layout featuring Govt. of Karnataka approval.",
            "School Registration Code: Dynamically binds ED 45 ACT registration details.",
            "Dynamic Student Records: Fills student name, grade, date of birth, and promotion details.",
            "Transfer Certificates (T.C.): Pupil Leaving Document featuring details table (name, parent, DOB).",
            "Single-Page Printing: Custom CSS rules avoiding print overlaps."
        ],
        "promotion_certificate.png"
    )

    # ==========================================
    # SLIDE 8: HOLIDAY PLANNER & MAILER
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8, COLOR_BG)
    add_slide_layout(
        slide8,
        "📅 School Holidays & Bulk Mailer",
        "Holiday Calendar & Mailer Portal",
        [
            "Published Calendar: Homepage calendar showing scheduled breaks year-on-year.",
            "Ordinal Date Format: Custom dates formatted in title-cased style (e.g., 1st Jan 2026, 19th Mar 2026).",
            "Interactive CRUD Panel: Full-width table allowing easy additions, edits, and deletions.",
            "Scheduled Auto-Mailer: Ticking 'Bulk Email Parents' checkbox sends auto-emails to parents.",
            "Custom Bulk Holiday Mailer: Specialized form to send weather alarms or ad-hoc school breaks."
        ],
        "dashboard_holidays.png"
    )

    # ==========================================
    # SLIDE 9: THANK YOU SLIDE
    # ==========================================
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9, COLOR_PRIMARY)

    left_accent9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
    left_accent9.fill.solid()
    left_accent9.fill.fore_color.rgb = COLOR_SECONDARY
    left_accent9.line.fill.background()

    thank_box = slide9.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.0))
    tf9 = thank_box.text_frame
    tf9.word_wrap = True
    p9 = tf9.paragraphs[0]
    p9.text = "Thank You!"
    p9.font.size = Pt(56)
    p9.font.bold = True
    p9.font.color.rgb = COLOR_WHITE
    p9.font.name = 'Helvetica'

    p9_sub = tf9.add_paragraph()
    p9_sub.text = "Kangaroo Kids Portal - Shaping Future Innovators\n\nFor questions or developer configuration, please access the Settings Panel."
    p9_sub.font.size = Pt(18)
    p9_sub.font.color.rgb = COLOR_SECONDARY
    p9_sub.font.name = 'Helvetica'
    p9_sub.space_before = Pt(20)

    # Save Presentation
    output_path = r"e:\AI_Applications\SchoolApplication\Kangaroo_Kids_Portal_Features.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    create_deck()
